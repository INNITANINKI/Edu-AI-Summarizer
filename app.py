import os
import json
import streamlit as st
import google.generativeai as genai
import tempfile
import time

from db import init_sqlite, save_faq
from transcript import fetch_transcript, extract_video_id
from rag_pipeline import summarize_text, generate_quiz, store_in_chroma, answer_question_with_mode

# ---------------- CONFIG ----------------
GENIE_MODEL = "gemini-2.5-flash"
GENIE_KEY = os.getenv("GEMINI_API_KEY", "YOUR_API_KEY_HERE")
if GENIE_KEY:
    genai.configure(api_key=GENIE_KEY)
else:
    st.warning("⚠️ GEMINI_API_KEY not set. Doubt solving will not work.")

# ---------------- STREAMLIT PAGE ----------------
st.set_page_config(page_title="YouTube & Document Summarizer", layout="centered")

# ---------------- DASHBOARD CSS ----------------
st.markdown("""
<style>
.main-title {
    font-size: 3rem;
    font-weight: 900;
    text-align: center;
    color: #ffffff;
    margin-bottom: 0.2rem;
}
.sub-title {
    text-align: center;
    font-size: 1.5rem;
    color: #ffffff;
    margin-bottom: 2rem;
}
.section-card {
    background: #ffffff; 
    color: #000000;
    border-radius: 18px;
    padding: 25px;
    margin: 20px 0;
    box-shadow: 0 10px 25px rgba(0,0,0,0.2);
}
.section-card.summary { border-left: 6px solid #888888; }
.section-card.answer { border-left: 6px solid #888888; }
.stButton>button {
    background-color: #f0f0f0;
    color: #000000;
    font-weight: 700;
    border-radius: 12px;
    padding: 14px 30px;
    font-size: 1.3rem;
    border: none;
    transition: all 0.3s ease;
}
.stButton>button:hover {
    background-color: #dddddd;
    transform: scale(1.05);
    box-shadow: 0 8px 20px rgba(0,0,0,0.3);
}
</style>
""", unsafe_allow_html=True)

# ---------------- HEADER ----------------
st.markdown('<div class="main-title">🎥📄 Edu AI Summarizer</div>', unsafe_allow_html=True)
st.markdown('<div class="sub-title">Summarize YouTube videos or documents, solve doubts & take quizzes!</div>', unsafe_allow_html=True)

# ---------------- SESSION STATE ----------------
for key in ["quiz", "user_answers", "submitted", "text_content", "doc_id", "video_id", "last_num_questions"]:
    if key not in st.session_state:
        st.session_state[key] = [] if key == "quiz" else {} if key == "user_answers" else False if key == "submitted" else "" if key in ["text_content", "doc_id", "video_id"] else 0

# ---------------- DATABASE ----------------
conn = init_sqlite()

# ---------------- SIDEBAR ----------------
st.sidebar.markdown("## 🚀 Learning Gateway")
st.sidebar.header("Input Options")
input_type = st.sidebar.radio("Select Input Type:", ["YouTube URL", "Upload Document"])

# ---------------- YOUTUBE URL ----------------
if input_type == "YouTube URL":
    youtube_url = st.sidebar.text_input("Enter YouTube video URL:")
    if st.sidebar.button("Fetch Transcript") or (st.session_state.video_id and st.session_state.video_id == extract_video_id(youtube_url)):
        video_id = extract_video_id(youtube_url)
        if not video_id:
            st.sidebar.error("❌ Invalid YouTube URL!")
        else:
            st.session_state.video_id = video_id
            # Unique doc_id for RAG: video_id + timestamp
            st.session_state.doc_id = f"{video_id}_{int(time.time())}"

            cursor = conn.cursor()
            cursor.execute("SELECT transcript FROM videos WHERE video_id=?", (video_id,))
            row = cursor.fetchone()
            if row and row[0]:
                st.session_state.text_content = row[0]
                st.success("✅ Loaded transcript from database.")
            else:
                with st.spinner("Fetching transcript..."):
                    transcript = fetch_transcript(video_id)
                if transcript:
                    cursor.execute(
                        "INSERT OR REPLACE INTO videos (video_id, title, transcript, summary, quiz, faqs) VALUES (?,?,?,?,?,?)",
                        (video_id, "Unknown", transcript, "", "[]", "[]")
                    )
                    conn.commit()
                    st.session_state.text_content = transcript
                    st.success("✅ Transcript fetched and saved.")

            if st.session_state.text_content:
                store_in_chroma(doc_id=st.session_state.doc_id, text=st.session_state.text_content)

# ---------------- DOCUMENT UPLOAD ----------------
elif input_type == "Upload Document":
    uploaded_file = st.sidebar.file_uploader("Upload PPTX, PDF, DOCX, TXT", type=['pptx','pdf','docx','txt'])
    if uploaded_file:

        def extract_text_from_file(file):
            extension = file.name.split('.')[-1].lower()
            text = ""
            try:
                with tempfile.NamedTemporaryFile(delete=False, suffix=f".{extension}") as tmp_file:
                    tmp_file.write(file.read())
                    tmp_path = tmp_file.name

                if extension == "pdf":
                    import pdfplumber
                    with pdfplumber.open(tmp_path) as pdf:
                        for page in pdf.pages:
                            text += (page.extract_text() or "") + "\n"

                elif extension in ["ppt", "pptx"]:
                    from pptx import Presentation
                    prs = Presentation(tmp_path)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                text += shape.text + "\n"

                elif extension in ["doc", "docx"]:
                    from docx import Document
                    doc = Document(tmp_path)
                    for para in doc.paragraphs:
                        if para.text.strip():
                            text += para.text + "\n"

                elif extension == "txt":
                    with open(tmp_path, "r", encoding="utf-8") as f:
                        text = f.read()
            except Exception as e:
                st.sidebar.error(f"❌ Failed to extract text: {e}")

            return text.strip()

        uploaded_file.seek(0)
        st.session_state.text_content = extract_text_from_file(uploaded_file)

        if st.session_state.text_content:
            st.success(f"✅ Extracted text from {uploaded_file.name}")
            # Unique doc_id for RAG: filename + timestamp
            st.session_state.doc_id = f"{uploaded_file.name}_{int(time.time())}"
            store_in_chroma(doc_id=st.session_state.doc_id, text=st.session_state.text_content)
        else:
            st.warning(f"⚠️ No readable text found in {uploaded_file.name}")

# ---------------- MODE SELECTION ----------------
if st.session_state.text_content:
    mode = st.sidebar.radio("Choose an action:", ["📄 Get Summary", "❓ Solve Doubt", "📝 Take Quiz"])

    # ---------------- SUMMARY ----------------
    if mode == "📄 Get Summary":
        if st.button("Generate Summary"):
            with st.spinner("Generating summary..."):
                summary = summarize_text(st.session_state.text_content)
            st.markdown(f'<div class="section-card summary"><h3>Summary</h3><p>{summary}</p></div>', unsafe_allow_html=True)

    # ---------------- DOUBT SOLVER ----------------
    elif mode == "❓ Solve Doubt":
        question = st.text_input("Enter your question:")
        mode_choice = st.radio("Select answer mode:", ["Based on provided context (RAG)", "General answer (AI)"])
        if question and st.button("Get Answer"):
            doc_id = st.session_state.doc_id
            selected_mode = "contextual" if "RAG" in mode_choice else "general"
            with st.spinner("Finding answer..."):
                answer = answer_question_with_mode(question, doc_id=doc_id, mode=selected_mode)
            st.markdown(f'<div class="section-card answer"><h3>Answer</h3><p>{answer}</p></div>', unsafe_allow_html=True)
            # Save FAQs only for YouTube using original video_id
            if selected_mode == "contextual" and input_type == "YouTube URL" and st.session_state.video_id:
                save_faq(conn, st.session_state.video_id, question, answer)

    # ---------------- QUIZ ----------------
    elif mode == "📝 Take Quiz":
        num_questions = st.selectbox("Number of questions:", [5, 10, 15, 20], index=0)
        if (not st.session_state.quiz) or (st.session_state.last_num_questions != num_questions):
            if st.button("Generate Quiz"):
                with st.spinner("Generating quiz..."):
                    st.session_state.quiz = generate_quiz(st.session_state.text_content, num_questions)
                    st.session_state.user_answers = {}
                    st.session_state.submitted = False
                    st.session_state.last_num_questions = num_questions

        if st.session_state.quiz and not st.session_state.submitted:
            for i, q in enumerate(st.session_state.quiz):
                if isinstance(q, dict) and "question" in q and "options" in q:
                    st.write(f"**Q{i+1}: {q['question']}**")
                    st.session_state.user_answers[i] = st.radio("Choose an option:", q["options"], key=f"q_{i}", index=None)
            if st.button("Submit Quiz"):
                if None in st.session_state.user_answers.values():
                    st.warning("⚠️ Please answer all questions before submitting.")
                else:
                    st.session_state.submitted = True

        if st.session_state.quiz and st.session_state.submitted:
            score = 0
            for i, q in enumerate(st.session_state.quiz):
                if isinstance(q, dict):
                    st.write(f"**Q{i+1}: {q['question']}**")
                    user_answer = st.session_state.user_answers.get(i)
                    st.write(f"Your answer: {user_answer}")
                    correct_option = next((opt for opt in q["options"] if opt.strip().upper().startswith(q["answer"].upper())), None)
                    if user_answer == correct_option:
                        st.success("✅ Correct!")
                        score += 1
                    else:
                        st.error(f"❌ Wrong! Correct answer: {correct_option}")
                    st.info(f"Explanation: {q.get('explanation','')}")
            st.subheader(f"🎯 Score: {score}/{len([q for q in st.session_state.quiz if isinstance(q, dict)])}")

else:
    st.info("⚠️ Enter a YouTube URL or upload a document first.")
